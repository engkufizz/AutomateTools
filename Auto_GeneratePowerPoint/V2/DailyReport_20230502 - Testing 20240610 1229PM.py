# Daily covid-19 report V2.0
# Uses pptx library to update the presentation with data from a text file.

import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import date, timedelta
import os

# Utility Functions
def loadValueFromText(file, lines=1):
    with open(file, 'r') as f:
        return f.readlines()[:lines]

def getMonthName(date_value, abbr=False):
    months = ["January", "February", "March", "April", "May", "June", 
              "July", "August", "September", "October", "November", "December"]
    if isinstance(date_value, str):
        month_int = int(date_value.split('-')[1])
    elif isinstance(date_value, int):
        month_int = date_value
    return months[month_int - 1][:3] if abbr else months[month_int - 1]

def getDateOfFewDaysAgo(date_str, days):
    d = date(*map(int, date_str.split('-')))
    delta = d - timedelta(days=days)
    return delta.year, delta.month, delta.day

if __name__ == '__main__':
    file_path = os.getcwd()
    file_name = "_Network KPI - Template_v3.pptx"
    data_source = "data.txt"

    file_full_path = os.path.join(file_path, file_name)
    source_data = loadValueFromText(os.path.join(file_path, data_source), 1)
    curr_data = source_data[0].split()
    file_date = curr_data[0].replace("-", "_", 2)
    file_report_name = f"Maxis_HSBA_OB_KPI_{file_date}.pptx"

    prs = Presentation(file_full_path)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for k, paragraph in enumerate(shape.text_frame.paragraphs):
                    text = paragraph.text

                    if i == 0:
                        curr_peak_percentage = f"{curr_data[4]}Gbps ({curr_data[5]})"
                        if j == 16 or j == 17:
                            paragraph.text = curr_peak_percentage
                            paragraph.font.name = 'Calibri Light'
                            paragraph.font.size = Pt(6.8)
                            paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                            if j == 17:
                                paragraph.font.bold = True
                                paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

                        if j == 18:
                            paragraph.text = curr_data[8] + "%"
                            paragraph.font.name = 'Calibri Light'
                            paragraph.font.size = Pt(6.8)
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                        if j == 14:
                            shape.rotation = 180 if curr_data[8].find("-") == 0 else 0

                        if j == 19:
                            paragraph.text = "decrease" if curr_data[8].find("-") == 0 else "increase"
                            paragraph.font.name = 'Calibri Light'
                            paragraph.font.size = Pt(6.8)
                            paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

                        if j == 20:
                            paragraph.text = curr_data[0]
                            paragraph.font.name = '微软雅黑'
                            paragraph.font.size = Pt(15)
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                    if i == 1:
                        if j == 2:
                            ttt = curr_data[0][8:10]
                            yesterday = str(int(ttt) - 1).zfill(2)
                            paragraph.text += f"{yesterday} {getMonthName(curr_data[0])} :"
                            paragraph.font.name = "Arial(Body)"
                            paragraph.font.size = Pt(16.5)
                            paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

                        if j == 5:
                            paragraph.text = curr_data[5]
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(22)
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

                        if j == 6:
                            paragraph.text = curr_data[4] + "G"
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(22)
                            paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

                        if j == 7:
                            if curr_data[8].find("-") == 0:
                                paragraph.text = "decrease"
                                curr_data[8] = curr_data[8][1:]
                            else:
                                paragraph.text = "increase"
                            paragraph.text += f" {curr_data[8]}%"
                            that_year, that_month, that_day = getDateOfFewDaysAgo(curr_data[0], 2)
                            paragraph.text += f" vs {str(that_day).zfill(2)} {getMonthName(that_month, False)}"
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(22)
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

                        if j == 8:
                            paragraph.text = ""

                        if j == 9:
                            if curr_data[9].find("-") == 0:
                                paragraph.text = "decrease "
                                curr_data[9] = curr_data[9][1:]
                            else:
                                paragraph.text = "increase "
                            last_week_year, last_week_month, last_week_day = getDateOfFewDaysAgo(curr_data[0], 7)
                            paragraph.text += f"{curr_data[9]} ({last_week_day}/{last_week_month}-{curr_data[7]}Gbps)"
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(22)
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

                        if j == 10:
                            paragraph.text = f"{curr_data[3]}Gbps - {curr_data[4]}Gbps (21:30 – 22:30)"
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(16.5)
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

                        if j == 11:
                            if float(curr_data[10].strip('%')) > float(curr_data[11].strip('%')):
                                curr_data[10], curr_data[11] = curr_data[11], curr_data[10]

                            paragraph.text = "increase" if curr_data[10].find("-") != 0 else "decrease"
                            curr_data[10] = curr_data[10].strip('-')
                            paragraph.text += f" {curr_data[10]}"
                            curr_data[11] = curr_data[11].strip('-')
                            paragraph.text += f" - {'decrease' if curr_data[11].find('-') == 0 else 'increase'} {curr_data[11]}"
                            paragraph.font.name = 'Arial (Headings)'
                            paragraph.font.size = Pt(16.5)
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

        if i == 1:
            img_path = os.path.join(file_path, "Traffic.png")
            left = Inches(0.3)
            top = Inches(1.67)
            heights = Inches(2.75)
            widths = Inches(9.43)
            slide.shapes.add_picture(img_path, left, top, height=heights, width=widths)

    prs.save(os.path.join(file_path, file_report_name))